"""
THISO Leasing Platform - User Guide PPT Generator
Generates a comprehensive PowerPoint presentation for all user roles
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

# ── Color Palette ─────────────────────────────────────────────────────────────
C_PRIMARY    = RGBColor(0x1E, 0x40, 0xAF)  # Blue-800
C_ACCENT     = RGBColor(0xF5, 0x9E, 0x0B)  # Amber-500
C_VIOLET     = RGBColor(0x76, 0x1D, 0xA8)  # Violet-700
C_GREEN      = RGBColor(0x05, 0x96, 0x69)  # Emerald-600
C_RED        = RGBColor(0xDC, 0x26, 0x26)  # Red-600
C_DARK       = RGBColor(0x11, 0x18, 0x27)  # Slate-900
C_GRAY       = RGBColor(0x64, 0x74, 0x8B)  # Slate-500
C_LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFF)  # Near white
C_DIVIDER    = RGBColor(0xE2, 0xE8, 0xF0)  # Slate-200
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_AMBER_LIGHT= RGBColor(0xFF, 0xF7, 0xED)
C_VIOLET_LIGHT=RGBColor(0xF5, 0xF3, 0xFF)

W  = Inches(13.33)   # Widescreen 16:9
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # Blank layout

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=12, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    try:
        run.font.name = 'Calibri'
    except Exception:
        pass
    return tb

def add_badge(slide, text, x, y, fill_color, text_color=C_WHITE, font_size=8.5, w=Inches(1.4), h=Inches(0.3)):
    r = add_rect(slide, x, y, w, h, fill=fill_color)
    add_text(slide, text, x, y + Pt(2), w, h, font_size=font_size, bold=True,
             color=text_color, align=PP_ALIGN.CENTER)
    return r

def heading_slide(title, subtitle=None, accent=C_PRIMARY):
    """Cover / section heading slide"""
    slide = prs.slides.add_slide(BLANK)

    # Background
    add_rect(slide, 0, 0, W, H, fill=accent)

    # White stripe
    add_rect(slide, 0, Inches(2.8), W, Inches(2.2), fill=C_WHITE)

    # Decorative circles
    add_rect(slide, Inches(-0.5), Inches(-0.5), Inches(2.5), Inches(2.5), fill=RGBColor(0xFF,0xFF,0xFF))
    circle = slide.shapes[-1]
    circle.fill.solid(); circle.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    circle.line.fill.background()
    # Note: rounded rect via XML would be complex; simple rectangle is fine

    add_text(slide, title,
             Inches(1.2), Inches(3.0), Inches(10.9), Inches(1.4),
             font_size=36, bold=True, color=accent, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.7),
                 font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    # THISO brand tag
    add_text(slide, 'THISO LEASING PLATFORM',
             Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
             font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    return slide

def content_slide(title, accent=C_PRIMARY):
    """Standard content slide with colored header bar"""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, W, Inches(1.1), fill=accent)
    add_text(slide, title, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.75),
             font_size=22, bold=True, color=C_WHITE)
    # Light background for content area
    add_rect(slide, 0, Inches(1.1), W, Inches(6.4), fill=C_LIGHT_BG)
    # Footer
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill=accent)
    add_text(slide, 'THISO Leasing Platform  |  Hướng dẫn người dùng  |  2026',
             Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.3),
             font_size=8, color=C_WHITE, align=PP_ALIGN.CENTER)
    return slide

def role_card(slide, x, y, w, h, role_name, email, password, modules, color):
    """A single role card with info"""
    add_rect(slide, x, y, w, h, fill=C_WHITE,
             line_color=color, line_width=Pt(1.5))
    # Color top bar
    add_rect(slide, x, y, w, Inches(0.42), fill=color)
    add_text(slide, role_name, x + Inches(0.12), y + Inches(0.05), w - Inches(0.24), Inches(0.35),
             font_size=10.5, bold=True, color=C_WHITE)
    # Email
    add_text(slide, '✉  ' + email, x + Inches(0.12), y + Inches(0.5), w - Inches(0.24), Inches(0.28),
             font_size=8.5, color=C_GRAY)
    add_text(slide, '🔑  ' + password, x + Inches(0.12), y + Inches(0.75), w - Inches(0.24), Inches(0.28),
             font_size=8.5, color=C_GRAY)
    # Modules
    add_text(slide, 'Phân hệ: ' + modules, x + Inches(0.12), y + Inches(1.0), w - Inches(0.24), Inches(h - Inches(1.12)),
             font_size=7.5, color=C_DARK, wrap=True)

def step_box(slide, x, y, num, title, desc, color=C_PRIMARY, w=Inches(2.3)):
    """Numbered step box for workflow diagrams"""
    bh = Inches(1.5)
    add_rect(slide, x, y, w, bh, fill=C_WHITE, line_color=color, line_width=Pt(1.5))
    # Number circle (simulated with small rect)
    add_rect(slide, x + Inches(0.12), y + Inches(0.1), Inches(0.42), Inches(0.42), fill=color)
    add_text(slide, str(num), x + Inches(0.12), y + Inches(0.1), Inches(0.42), Inches(0.42),
             font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.62), y + Inches(0.12), w - Inches(0.74), Inches(0.42),
             font_size=9.5, bold=True, color=color)
    add_text(slide, desc, x + Inches(0.12), y + Inches(0.58), w - Inches(0.24), Inches(0.85),
             font_size=8, color=C_GRAY, wrap=True)

def arrow(slide, x1, y, x2):
    """Simple horizontal arrow"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    connector = slide.shapes.add_connector(1, x1, y, x2, y)
    connector.line.color.rgb = C_GRAY
    connector.line.width = Pt(1.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# Deep blue full background
add_rect(slide, 0, 0, W, H, fill=C_PRIMARY)
# White content panel
add_rect(slide, Inches(1.0), Inches(1.5), Inches(11.33), Inches(4.5), fill=C_WHITE)
# Accent stripe
add_rect(slide, Inches(1.0), Inches(1.5), Inches(0.18), Inches(4.5), fill=C_ACCENT)

add_text(slide, 'THISO LEASING PLATFORM',
         Inches(1.5), Inches(1.8), Inches(10.5), Inches(0.6),
         font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
add_text(slide, 'Hướng dẫn\nSử dụng Hệ thống',
         Inches(1.5), Inches(2.3), Inches(10.5), Inches(2.2),
         font_size=42, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)
add_text(slide, 'Dành cho tất cả người dùng theo vai trò & phân quyền',
         Inches(1.5), Inches(4.3), Inches(10.5), Inches(0.5),
         font_size=14, color=C_GRAY, align=PP_ALIGN.LEFT)
add_text(slide, 'Phiên bản 1.0  |  Tháng 6/2026  |  THISO Mall Sala',
         Inches(1.5), Inches(5.7), Inches(10.5), Inches(0.35),
         font_size=10, color=C_GRAY, align=PP_ALIGN.LEFT)
# Bottom bar
add_rect(slide, 0, Inches(6.9), W, Inches(0.6), fill=C_ACCENT)
add_text(slide, 'Tài liệu này được tạo bởi THISO Digital Team  —  Bảo mật nội bộ',
         Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4),
         font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MỤC LỤC
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Mục lục', accent=C_PRIMARY)

items = [
    ('01', 'Tổng quan hệ thống',        'Giới thiệu THISO Leasing Platform, kiến trúc phân hệ'),
    ('02', 'Cơ cấu tổ chức & Tài khoản','Danh sách vai trò, email đăng nhập, mật khẩu'),
    ('03', 'Quy trình Cho thuê Dài hạn','CRM Lead → Booking → Proposal → Hợp đồng → Tenant'),
    ('04', 'Quy trình Cho thuê Ngắn hạn','Đặt Slot → Xác nhận → Thanh toán → Hoàn thành'),
    ('05', 'Quản lý Vận hành',           'Giám sát mặt bằng, thi công Fitout, xử lý sự cố'),
    ('06', 'Kế toán & Tài chính',        'Hóa đơn, thu tiền, báo cáo doanh thu, AR Aging'),
    ('07', 'Dashboard & Báo cáo',        'Dashboard CEO, phân tích đa trung tâm'),
    ('08', 'Phân quyền chi tiết',        'Ma trận quyền truy cập theo vai trò'),
]
for i, (num, title, desc) in enumerate(items):
    row = i // 2
    col = i %  2
    x = Inches(0.5) + col * Inches(6.4)
    y = Inches(1.3) + row * Inches(1.4)
    add_rect(slide, x, y, Inches(6.0), Inches(1.25), fill=C_WHITE,
             line_color=C_DIVIDER, line_width=Pt(1))
    add_rect(slide, x, y, Inches(0.65), Inches(1.25), fill=C_PRIMARY)
    add_text(slide, num, x, y, Inches(0.65), Inches(1.25),
             font_size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.75), y + Inches(0.15), Inches(5.1), Inches(0.45),
             font_size=13, bold=True, color=C_DARK)
    add_text(slide, desc, x + Inches(0.75), y + Inches(0.58), Inches(5.1), Inches(0.55),
             font_size=9, color=C_GRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TỔNG QUAN HỆ THỐNG
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('01. Tổng quan Hệ thống THISO Leasing Platform')

modules = [
    ('🏬', 'Mặt bằng',   'Quản lý lô thuê, tầng, khu vực, bản đồ số', C_PRIMARY),
    ('👥', 'CRM',        'Lead, khách hàng, pipeline đàm phán', C_GREEN),
    ('📌', 'Booking DH', 'Giữ lô dài hạn, hàng đợi ưu tiên', RGBColor(0xD9,0x77,0x06)),
    ('🗓', 'Booking NH', 'Slot pop-up/kiosk theo ngày/giờ/tháng', C_VIOLET),
    ('📋', 'Proposal',   'Đề xuất giá, phê duyệt đa cấp', C_RED),
    ('📄', 'Hợp đồng',   'Ký kết, gia hạn, chấm dứt hợp đồng', C_PRIMARY),
    ('🔧', 'Vận hành',   'Fitout, bảo trì, ticket sự cố', C_GREEN),
    ('💰', 'Kế toán',    'Hóa đơn, thu tiền, AR Aging, báo cáo', RGBColor(0xD9,0x77,0x06)),
    ('📊', 'Dashboard',  'KPI tổng hợp, phân tích đa trung tâm', C_PRIMARY),
]
cols = 3
for i, (icon, name, desc, color) in enumerate(modules):
    col = i % cols
    row = i // cols
    x = Inches(0.35) + col * Inches(4.3)
    y = Inches(1.3) + row * Inches(1.9)
    add_rect(slide, x, y, Inches(3.9), Inches(1.75), fill=C_WHITE,
             line_color=color, line_width=Pt(1.5))
    add_rect(slide, x, y, Inches(0.55), Inches(1.75), fill=color)
    add_text(slide, icon, x, y + Inches(0.55), Inches(0.55), Inches(0.65),
             font_size=18, align=PP_ALIGN.CENTER)
    add_text(slide, name, x + Inches(0.65), y + Inches(0.2), Inches(3.1), Inches(0.4),
             font_size=12, bold=True, color=color)
    add_text(slide, desc, x + Inches(0.65), y + Inches(0.6), Inches(3.1), Inches(1.0),
             font_size=8.5, color=C_GRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CƠ CẤU TỔ CHỨC & TÀI KHOẢN
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('02. Cơ cấu Tổ chức & Tài khoản Hệ thống')

accounts = [
    ('Tổng Giám Đốc (CEO)', 'ceo@thiso.com', 'User123!',
     'Dashboard, CRM, Proposals, Reports, Analytics, Cross-mall', C_DARK),
    ('Giám Đốc Mall', 'do.thi.nga@thiso.com', 'Thiso@2024',
     'Dashboard, Mặt bằng, CRM, Booking, Proposals, Hợp đồng, Vận hành', C_PRIMARY),
    ('Trưởng Nhóm Sale DH', 'nguyen.thu.hoa@thiso.com', 'Thiso@2024',
     'CRM, Booking dài hạn, Proposal, Hợp đồng, Báo cáo', RGBColor(0xD9,0x77,0x06)),
    ('Sale Executive DH', 'le.minh.tuan@thiso.com', 'Thiso@2024',
     'CRM Leads, Booking dài hạn, Mặt bằng', RGBColor(0xD9,0x77,0x06)),
    ('Sale Executive DH', 'pham.thi.lan@thiso.com', 'Thiso@2024',
     'CRM Leads, Booking dài hạn, Mặt bằng', RGBColor(0xD9,0x77,0x06)),
    ('Trưởng Nhóm Sale NH', 'tran.van.duc@thiso.com', 'Thiso@2024',
     'Booking ngắn hạn, Slot, CRM, Báo cáo', C_VIOLET),
    ('Sale Executive NH', 'hoang.thi.mai@thiso.com', 'Thiso@2024',
     'Booking ngắn hạn, Slot pop-up, CRM', C_VIOLET),
    ('Sale Executive NH', 'vu.quoc.binh@thiso.com', 'Thiso@2024',
     'Booking ngắn hạn, Slot pop-up, CRM', C_VIOLET),
    ('Vận Hành', 'trinh.van.long@thiso.com', 'Thiso@2024',
     'Ticket, Fitout, Bảo trì, Thông báo', C_GREEN),
    ('Kế Toán Trưởng', 'bui.thi.hoa@thiso.com', 'Thiso@2024',
     'Billing, Hóa đơn, AR Aging, Báo cáo, SAP sync', RGBColor(0x0F,0x76,0x6E)),
    ('Kế Toán Viên', 'dang.van.quang@thiso.com', 'Thiso@2024',
     'Billing, Hóa đơn, Báo cáo', RGBColor(0x0F,0x76,0x6E)),
]

cols = 3
card_w = Inches(4.1)
card_h = Inches(1.6)
for i, (role, email, pwd, mods, color) in enumerate(accounts):
    col = i % cols
    row = i // cols
    x = Inches(0.3) + col * (card_w + Inches(0.18))
    y = Inches(1.25) + row * (card_h + Inches(0.12))
    if row > 2:
        continue  # Max 3 rows visible
    add_rect(slide, x, y, card_w, card_h, fill=C_WHITE,
             line_color=color, line_width=Pt(1.2))
    add_rect(slide, x, y, card_w, Inches(0.38), fill=color)
    add_text(slide, role, x + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2), Inches(0.32),
             font_size=9.5, bold=True, color=C_WHITE)
    add_text(slide, '✉ ' + email, x + Inches(0.1), y + Inches(0.44), card_w - Inches(0.2), Inches(0.26),
             font_size=7.5, color=C_GRAY)
    add_text(slide, '🔑 ' + pwd, x + Inches(0.1), y + Inches(0.68), card_w - Inches(0.2), Inches(0.26),
             font_size=7.5, color=C_GRAY)
    add_text(slide, mods, x + Inches(0.1), y + Inches(0.94), card_w - Inches(0.2), Inches(0.55),
             font_size=7, color=C_DARK, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HƯỚNG DẪN ĐĂNG NHẬP
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Đăng nhập hệ thống', accent=C_PRIMARY)

add_text(slide, 'Truy cập:', Inches(0.5), Inches(1.25), Inches(3), Inches(0.4),
         font_size=12, bold=True, color=C_DARK)
add_rect(slide, Inches(0.5), Inches(1.65), Inches(5.5), Inches(0.5), fill=C_PRIMARY)
add_text(slide, 'http://localhost:8080  (nội bộ)',
         Inches(0.6), Inches(1.68), Inches(5.3), Inches(0.45),
         font_size=13, bold=True, color=C_WHITE)

steps_login = [
    ('1', 'Mở trình duyệt', 'Chrome hoặc Edge, truy cập địa chỉ hệ thống'),
    ('2', 'Nhập thông tin', 'Email và mật khẩu theo danh sách tài khoản'),
    ('3', 'Chọn Mall', 'Sau đăng nhập, chọn THISO Mall Sala từ menu trên'),
    ('4', 'Vào Dashboard', 'Hệ thống hiển thị theo vai trò của bạn'),
]
for i, (num, title, desc) in enumerate(steps_login):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(2.5)
    step_box(slide, x, y, int(num), title, desc, color=C_PRIMARY, w=Inches(2.85))

# Tips box
add_rect(slide, Inches(0.5), Inches(4.3), Inches(12.33), Inches(2.85),
         fill=RGBColor(0xEF,0xF6,0xFF), line_color=C_PRIMARY, line_width=Pt(1))
add_text(slide, '💡  Lưu ý quan trọng', Inches(0.7), Inches(4.45), Inches(11.9), Inches(0.4),
         font_size=12, bold=True, color=C_PRIMARY)
tips = [
    '• Mật khẩu phân biệt HOA / thường — nhập chính xác như trong tài liệu này',
    '• Sau lần đăng nhập đầu tiên, khuyến khích đổi mật khẩu tại menu hồ sơ cá nhân',
    '• Token đăng nhập có hiệu lực 7 ngày — không cần đăng nhập lại hàng ngày',
    '• Nếu quên mật khẩu: liên hệ Admin hệ thống (admin@thiso.com) để reset',
    '• Mỗi tài khoản chỉ được truy cập vào các phân hệ phù hợp với vai trò',
]
add_text(slide, '\n'.join(tips), Inches(0.7), Inches(4.85), Inches(11.9), Inches(2.2),
         font_size=10, color=C_DARK, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SECTION: SALE DÀI HẠN
# ═══════════════════════════════════════════════════════════════════════════════
heading_slide('03. Quy trình Cho thuê Dài hạn',
              'Dành cho Sale Executive & Leasing Manager — Team Dài Hạn',
              accent=RGBColor(0xD9,0x77,0x06))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WORKFLOW DÀI HẠN OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Tổng quan quy trình Cho thuê Dài hạn', accent=RGBColor(0xD9,0x77,0x06))

stages = [
    ('Lead\n(CRM)', 'Tiếp nhận khách\ntìm mặt bằng', C_GREEN),
    ('Booking\n(Giữ lô)', 'Đặt cọc chỗ\nvào hàng đợi', RGBColor(0xD9,0x77,0x06)),
    ('Proposal\n(Đề xuất)', 'Lập đề xuất\ngiá chi tiết', C_VIOLET),
    ('Phê duyệt\n(Approval)', 'Manager/Director\nký duyệt', C_RED),
    ('Hợp đồng\n(Contract)', 'Ký kết hợp\nđồng thuê', C_PRIMARY),
    ('Tenant\n(Vận hành)', 'Bàn giao, thi\ncông, khai thác', C_GREEN),
]
stage_w = Inches(1.85)
gap = Inches(0.2)
start_x = Inches(0.35)
y = Inches(1.7)

for i, (label, desc, color) in enumerate(stages):
    x = start_x + i * (stage_w + gap)
    # Arrow before (except first)
    if i > 0:
        ax = x - gap
        add_rect(slide, ax - Inches(0.05), y + Inches(0.55), Inches(0.28), Inches(0.1), fill=C_GRAY)

    add_rect(slide, x, y, stage_w, Inches(1.3), fill=color)
    add_text(slide, label, x, y + Inches(0.15), stage_w, Inches(0.7),
             font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, y + Inches(1.3), stage_w, Inches(0.9), fill=C_WHITE,
             line_color=color, line_width=Pt(1))
    add_text(slide, desc, x + Inches(0.08), y + Inches(1.35), stage_w - Inches(0.16), Inches(0.8),
             font_size=8, color=C_DARK, wrap=True, align=PP_ALIGN.CENTER)

# Role assignment
add_text(slide, 'Vai trò thực hiện:', Inches(0.35), Inches(3.35), Inches(3), Inches(0.4),
         font_size=11, bold=True, color=C_DARK)
roles_map = [
    ('Sale Executive', 'Tạo Lead, tạo Booking', RGBColor(0xD9,0x77,0x06)),
    ('Leasing Manager', 'Convert Proposal, submit duyệt', C_PRIMARY),
    ('Mall Director', 'Phê duyệt Proposal', C_RED),
    ('Finance', 'Xem hợp đồng, lập hóa đơn', RGBColor(0x0F,0x76,0x6E)),
]
for i, (role, task, color) in enumerate(roles_map):
    x = Inches(0.35) + i * Inches(3.2)
    add_rect(slide, x, Inches(3.75), Inches(2.9), Inches(0.55), fill=C_WHITE,
             line_color=color, line_width=Pt(1.2))
    add_rect(slide, x, Inches(3.75), Inches(0.18), Inches(0.55), fill=color)
    add_text(slide, role, x + Inches(0.28), Inches(3.78), Inches(2.5), Inches(0.25),
             font_size=9.5, bold=True, color=color)
    add_text(slide, task, x + Inches(0.28), Inches(4.02), Inches(2.5), Inches(0.22),
             font_size=7.5, color=C_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BƯỚC 1: TẠO LEAD TRONG CRM
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Bước 1 — Tạo Lead trong CRM (Sale Executive)', accent=RGBColor(0xD9,0x77,0x06))

steps = [
    ('1', 'Vào menu CRM', 'Sidebar trái → "CRM" → Chọn "Leads"'),
    ('2', 'Bấm "+ Tạo Lead"', 'Nút xanh góc trên phải màn hình'),
    ('3', 'Điền thông tin', 'Tên thương hiệu, công ty, liên hệ, nguồn, loại mặt bằng mong muốn'),
    ('4', 'Thiết lập ngân sách', 'Diện tích mong muốn (m²), giá thuê tối đa (₫/m²)'),
    ('5', 'Gán cho mình', 'Mục "Phụ trách" → chọn tên của bạn'),
    ('6', 'Lưu Lead', 'Bấm "Tạo Lead" — Lead xuất hiện ở cột NEW trong Pipeline'),
]
for i, (num, title, desc) in enumerate(steps):
    row = i // 3
    col = i % 3
    x = Inches(0.35) + col * Inches(4.3)
    y = Inches(1.35) + row * Inches(2.1)
    add_rect(slide, x, y, Inches(3.95), Inches(1.85), fill=C_WHITE,
             line_color=RGBColor(0xD9,0x77,0x06), line_width=Pt(1.2))
    add_rect(slide, x, y, Inches(0.55), Inches(0.55), fill=RGBColor(0xD9,0x77,0x06))
    add_text(slide, num, x, y, Inches(0.55), Inches(0.55),
             font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.65), y + Inches(0.1), Inches(3.18), Inches(0.38),
             font_size=11, bold=True, color=RGBColor(0xD9,0x77,0x06))
    add_text(slide, desc, x + Inches(0.12), y + Inches(0.6), Inches(3.7), Inches(1.1),
             font_size=9.5, color=C_DARK, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — BƯỚC 2: TẠO BOOKING DÀI HẠN
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Bước 2 — Tạo Booking Dài hạn (Giữ lô)', accent=RGBColor(0xD9,0x77,0x06))

add_text(slide, 'Có 2 cách tạo Booking:', Inches(0.4), Inches(1.2), Inches(12), Inches(0.4),
         font_size=12, bold=True, color=C_DARK)

# Method A
add_rect(slide, Inches(0.4), Inches(1.65), Inches(5.9), Inches(3.5), fill=C_AMBER_LIGHT,
         line_color=RGBColor(0xD9,0x77,0x06), line_width=Pt(1.5))
add_text(slide, 'Cách A — Từ trang Mặt bằng', Inches(0.55), Inches(1.75), Inches(5.6), Inches(0.45),
         font_size=12, bold=True, color=RGBColor(0xD9,0x77,0x06))
steps_a = [
    '① Vào menu "Mặt bằng" → chọn Mall',
    '② Bấm vào lô thuê (unit) muốn giữ',
    '③ Tab "Booking" → bấm "Tạo booking"',
    '④ Chọn Lead từ CRM, điền yêu cầu',
    '⑤ Xác nhận — Unit chuyển sang màu vàng (BOOKING)',
]
for i, s in enumerate(steps_a):
    add_text(slide, s, Inches(0.6), Inches(2.2) + i * Inches(0.5), Inches(5.6), Inches(0.45),
             font_size=10, color=C_DARK)

# Method B
add_rect(slide, Inches(6.9), Inches(1.65), Inches(5.9), Inches(3.5), fill=C_AMBER_LIGHT,
         line_color=RGBColor(0xD9,0x77,0x06), line_width=Pt(1.5))
add_text(slide, 'Cách B — Từ trang Booking', Inches(7.05), Inches(1.75), Inches(5.6), Inches(0.45),
         font_size=12, bold=True, color=RGBColor(0xD9,0x77,0x06))
steps_b = [
    '① Vào menu "Booking" → Tab "Giữ lô dài hạn"',
    '② Bấm nút "Tạo booking" (màu cam)',
    '③ Chọn Unit → chọn Lead CRM',
    '④ Điền diện tích, thời hạn, giá kỳ vọng',
    '⑤ Hệ thống tự xếp vào hàng đợi ưu tiên',
]
for i, s in enumerate(steps_b):
    add_text(slide, s, Inches(7.05), Inches(2.2) + i * Inches(0.5), Inches(5.6), Inches(0.45),
             font_size=10, color=C_DARK)

# Note about priority queue
add_rect(slide, Inches(0.4), Inches(5.3), Inches(12.4), Inches(0.9), fill=C_WHITE,
         line_color=C_PRIMARY, line_width=Pt(1))
add_text(slide,
         '📌  Hàng đợi ưu tiên: Một unit có thể có nhiều booking — khách ưu tiên #1 được giữ slot, '
         'các booking #2, #3 ở trạng thái PENDING chờ. Nếu #1 hủy, hệ thống tự động kích hoạt #2.',
         Inches(0.55), Inches(5.38), Inches(12.1), Inches(0.75),
         font_size=9.5, color=C_DARK, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — BƯỚC 3: TẠO PROPOSAL
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Bước 3 — Lập Đề xuất (Proposal) — Leasing Manager', accent=RGBColor(0xD9,0x77,0x06))

steps3 = [
    ('1', 'Vào trang Booking', 'Tab "Giữ lô dài hạn" → click vào booking của khách'),
    ('2', 'Bấm "Lập Đề xuất"', 'Nút xanh lá trong detail panel bên phải'),
    ('3', 'Điền thông tin giá', 'Diện tích, thời hạn thuê, giá/m², phí CAM, tiền đặt cọc'),
    ('4', 'Điều khoản thêm', 'Rent-free (ngày), Escalation (%/năm), ghi chú đặc biệt'),
    ('5', 'Xem preview', 'Hệ thống tính tổng giá thuê ước tính tự động'),
    ('6', 'Submit để duyệt', 'Bấm "Gửi duyệt" — Proposal vào hàng chờ phê duyệt'),
]
for i, (num, title, desc) in enumerate(steps3):
    row = i // 3; col = i % 3
    x = Inches(0.35) + col * Inches(4.3)
    y = Inches(1.35) + row * Inches(2.1)
    add_rect(slide, x, y, Inches(3.95), Inches(1.85), fill=C_WHITE,
             line_color=RGBColor(0xD9,0x77,0x06), line_width=Pt(1.2))
    add_rect(slide, x, y, Inches(0.55), Inches(0.55), fill=C_VIOLET)
    add_text(slide, num, x, y, Inches(0.55), Inches(0.55),
             font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.65), y + Inches(0.1), Inches(3.18), Inches(0.38),
             font_size=11, bold=True, color=C_VIOLET)
    add_text(slide, desc, x + Inches(0.12), y + Inches(0.6), Inches(3.7), Inches(1.1),
             font_size=9.5, color=C_DARK, wrap=True)

add_rect(slide, Inches(0.35), Inches(5.55), Inches(12.62), Inches(0.75), fill=C_VIOLET_LIGHT,
         line_color=C_VIOLET, line_width=Pt(1))
add_text(slide, '⚡  Quy tắc phê duyệt tự động: Proposal thường → Manager → Director. '
         'Nếu chiết khấu >10% hoặc giá dưới sàn → cần CEO phê duyệt thêm.',
         Inches(0.55), Inches(5.62), Inches(12.2), Inches(0.6), font_size=9.5, color=C_VIOLET, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SECTION: SALE NGẮN HẠN
# ═══════════════════════════════════════════════════════════════════════════════
heading_slide('04. Quy trình Cho thuê Ngắn hạn (Slot)',
              'Dành cho Sale Executive & Manager — Team Ngắn Hạn',
              accent=C_VIOLET)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SLOT BOOKING OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Tổng quan Booking Slot Ngắn hạn', accent=C_VIOLET)

# 3 types
types = [
    ('📅', 'Theo Ngày\n(DAILY)', 'Thuê theo số ngày\nPop-up 1-30 ngày\nGiá: X₫/m²/ngày', RGBColor(0x0E,0xA5,0xE9)),
    ('⏰', 'Theo Giờ\n(HOURLY)', 'Thuê theo giờ\nSự kiện, demo\nGiá: X₫/giờ', RGBColor(0xF9,0x73,0x16)),
    ('🗓', 'Theo Tháng\n(MONTHLY)', 'Thuê cố định 1+ tháng\nQuầy hàng thường xuyên\nGiá: X₫/m²/tháng', RGBColor(0x14,0xB8,0xA6)),
]
for i, (icon, label, desc, color) in enumerate(types):
    x = Inches(0.5) + i * Inches(4.1)
    add_rect(slide, x, Inches(1.3), Inches(3.7), Inches(3.2), fill=C_WHITE,
             line_color=color, line_width=Pt(2))
    add_rect(slide, x, Inches(1.3), Inches(3.7), Inches(0.95), fill=color)
    add_text(slide, icon + '  ' + label, x, Inches(1.42), Inches(3.7), Inches(0.72),
             font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + Inches(0.2), Inches(2.35), Inches(3.3), Inches(1.95),
             font_size=11, color=C_DARK, wrap=True, align=PP_ALIGN.CENTER)

# Flow
add_text(slide, 'Quy trình:', Inches(0.5), Inches(4.7), Inches(2.5), Inches(0.4),
         font_size=12, bold=True, color=C_DARK)
flow = [
    ('Tạo Booking', 'Sale Executive\ntạo đặt slot', C_VIOLET),
    ('Chờ XN', 'Trạng thái\nPENDING', RGBColor(0x60,0xA5,0xFA)),
    ('Xác nhận', 'Manager\nbấm Confirm', C_GREEN),
    ('Hoàn thành', 'Sau ngày kết\nthúc → DONE', RGBColor(0x10,0xB9,0x81)),
]
for i, (label, desc, color) in enumerate(flow):
    x = Inches(0.5) + i * Inches(3.1)
    add_rect(slide, x, Inches(5.2), Inches(2.7), Inches(1.2), fill=color)
    add_text(slide, label, x, Inches(5.25), Inches(2.7), Inches(0.4),
             font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x, Inches(5.62), Inches(2.7), Inches(0.65),
             font_size=9, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, '→', x + Inches(2.7), Inches(5.55), Inches(0.4), Inches(0.35),
                 font_size=16, bold=True, color=C_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TẠO SLOT BOOKING
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Tạo Booking Slot từ trang Booking', accent=C_VIOLET)

steps_slot = [
    ('1', 'Vào menu Booking', 'Menu "Booking" → Tab tím "Đặt slot ngắn hạn"'),
    ('2', 'Bấm "+ Tạo booking slot"', 'Nút tím góc trên trái bảng danh sách'),
    ('3', 'Chọn Lô + Slot', 'Dropdown Lô thuê → Dropdown Slot (cascade tự động)'),
    ('4', 'Chọn Khách hàng', 'Toggle "Lead CRM" hoặc "Customer profile" → search tên'),
    ('5', 'Chọn loại + Thời gian', 'DAILY/HOURLY/MONTHLY → Ngày bắt đầu → Ngày kết thúc'),
    ('6', 'Xem giá & xác nhận', 'Giá tự tính, kiểm tra → Bấm "Tạo booking"'),
]
for i, (num, title, desc) in enumerate(steps_slot):
    row = i // 3; col = i % 3
    x = Inches(0.35) + col * Inches(4.3)
    y = Inches(1.35) + row * Inches(2.1)
    add_rect(slide, x, y, Inches(3.95), Inches(1.85), fill=C_VIOLET_LIGHT,
             line_color=C_VIOLET, line_width=Pt(1.2))
    add_rect(slide, x, y, Inches(0.55), Inches(0.55), fill=C_VIOLET)
    add_text(slide, num, x, y, Inches(0.55), Inches(0.55),
             font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + Inches(0.65), y + Inches(0.1), Inches(3.18), Inches(0.38),
             font_size=11, bold=True, color=C_VIOLET)
    add_text(slide, desc, x + Inches(0.12), y + Inches(0.6), Inches(3.7), Inches(1.1),
             font_size=9.5, color=C_DARK, wrap=True)

add_rect(slide, Inches(0.35), Inches(5.55), Inches(12.62), Inches(0.75), fill=C_WHITE,
         line_color=C_GREEN, line_width=Pt(1))
add_text(slide, '✅  Sau khi Manager bấm "Xác nhận": Slot booking chuyển sang CONFIRMED '
         '— slot đó bị khóa trong khoảng thời gian đã đặt, không ai đặt trùng được.',
         Inches(0.55), Inches(5.62), Inches(12.2), Inches(0.6), font_size=9.5, color=C_GREEN, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — SECTION: VẬN HÀNH
# ═══════════════════════════════════════════════════════════════════════════════
heading_slide('05. Quản lý Vận hành',
              'Dành cho Mall Director, Operation Staff — Team Vận Hành',
              accent=C_GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — VẬN HÀNH CHI TIẾT
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Phân hệ Quản lý Vận hành', accent=C_GREEN)

ops_modules = [
    ('🗺', 'Bản đồ Mặt bằng',
     '• Xem sơ đồ tầng, trạng thái lô\n• Màu sắc: Trống/Booking/Thuê/Thi công\n• Theo dõi realtime',
     C_PRIMARY),
    ('🔧', 'Ticket Sự cố',
     '• Tạo ticket HVAC, điện, nước, PCCC\n• Gán nhân viên xử lý\n• Theo dõi tiến độ, SLA\n• Lưu lịch sử',
     C_GREEN),
    ('🏗', 'Fitout Thi công',
     '• Quản lý dự án thi công của tenant\n• Checklist kiểm tra từng giai đoạn\n• Upload tài liệu thiết kế\n• Theo dõi mốc thời gian',
     C_RED),
    ('📢', 'Thông báo Mall',
     '• Đăng thông báo nội bộ\n• Gửi tới tenant, nhân viên\n• Theo lịch hoặc tức thì',
     RGBColor(0xD9,0x77,0x06)),
]
for i, (icon, title, content, color) in enumerate(ops_modules):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.45)
    y = Inches(1.3) + row * Inches(2.85)
    add_rect(slide, x, y, Inches(6.1), Inches(2.65), fill=C_WHITE,
             line_color=color, line_width=Pt(1.5))
    add_rect(slide, x, y, Inches(6.1), Inches(0.55), fill=color)
    add_text(slide, icon + '  ' + title, x + Inches(0.12), y + Inches(0.1),
             Inches(5.86), Inches(0.42), font_size=13, bold=True, color=C_WHITE)
    add_text(slide, content, x + Inches(0.15), y + Inches(0.65),
             Inches(5.8), Inches(1.85), font_size=9.5, color=C_DARK, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — SECTION: KẾ TOÁN
# ═══════════════════════════════════════════════════════════════════════════════
heading_slide('06. Kế toán & Tài chính',
              'Dành cho Finance Manager & Kế toán viên — Team Kế Toán',
              accent=RGBColor(0x0F,0x76,0x6E))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — KẾ TOÁN CHI TIẾT
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Phân hệ Kế toán & Tài chính', accent=RGBColor(0x0F,0x76,0x6E))

finance_modules = [
    ('💳', 'Hóa đơn (Invoice)',
     '• Xem danh sách hóa đơn theo tenant\n• Phát hành hóa đơn hàng tháng\n• Ghi nhận thanh toán\n• Export Excel/PDF',
     RGBColor(0x0F,0x76,0x6E)),
    ('📊', 'AR Aging',
     '• Theo dõi công nợ theo nhóm: <30/30-60/60-90/>90 ngày\n• Xác định tenant quá hạn\n• Gửi nhắc nợ tự động',
     C_RED),
    ('📈', 'Báo cáo Doanh thu',
     '• Revenue theo tháng, quý, năm\n• So sánh theo mall, theo danh mục\n• KPI thu tiền\n• Tích hợp SAP',
     C_PRIMARY),
    ('⏰', 'Lịch thanh toán',
     '• Lịch thu tiền tự động từ hợp đồng\n• Cảnh báo hóa đơn đến hạn\n• Phạt chậm trả tự động\n• Đối soát ngân hàng',
     RGBColor(0xD9,0x77,0x06)),
]
for i, (icon, title, content, color) in enumerate(finance_modules):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.45)
    y = Inches(1.3) + row * Inches(2.85)
    add_rect(slide, x, y, Inches(6.1), Inches(2.65), fill=C_WHITE,
             line_color=color, line_width=Pt(1.5))
    add_rect(slide, x, y, Inches(6.1), Inches(0.55), fill=color)
    add_text(slide, icon + '  ' + title, x + Inches(0.12), y + Inches(0.1),
             Inches(5.86), Inches(0.42), font_size=13, bold=True, color=C_WHITE)
    add_text(slide, content, x + Inches(0.15), y + Inches(0.65),
             Inches(5.8), Inches(1.85), font_size=9.5, color=C_DARK, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CEO DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════════
heading_slide('07. Dashboard & Báo cáo CEO',
              'Dành cho Tổng Giám Đốc & Ban Giám Đốc',
              accent=C_DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CEO MODULES
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Dashboard & Tính năng dành riêng CEO', accent=C_DARK)

ceo_items = [
    ('📊', 'Dashboard Tổng hợp',
     '• Occupancy Rate toàn mall realtime\n• Doanh thu tháng vs target\n• KPI leasing pipeline\n• Số lô trống, đang đàm phán, đã ký',
     C_PRIMARY),
    ('🌐', 'Cross-Mall Dashboard',
     '• So sánh chỉ số giữa nhiều mall\n• Tổng doanh thu toàn hệ thống\n• Ranking mall theo performance\n• Dữ liệu ngành hàng xuyên mall',
     C_VIOLET),
    ('✅', 'Phê duyệt Đặc biệt',
     '• Proposal có chiết khấu > 10%\n• Giá dưới sàn quy định\n• Các trường hợp exception\n• Lịch sử quyết định',
     C_RED),
    ('🤖', 'AI Assistant',
     '• Hỏi đáp về dữ liệu hệ thống\n• Phân tích sơ đồ mặt bằng bằng AI\n• Gợi ý cơ cấu ngành hàng\n• Dự báo occupancy',
     C_GREEN),
]
for i, (icon, title, content, color) in enumerate(ceo_items):
    col = i % 2; row = i // 2
    x = Inches(0.35) + col * Inches(6.45)
    y = Inches(1.3) + row * Inches(2.85)
    add_rect(slide, x, y, Inches(6.1), Inches(2.65), fill=C_WHITE,
             line_color=color, line_width=Pt(1.5))
    add_rect(slide, x, y, Inches(6.1), Inches(0.55), fill=color)
    add_text(slide, icon + '  ' + title, x + Inches(0.12), y + Inches(0.1),
             Inches(5.86), Inches(0.42), font_size=13, bold=True, color=C_WHITE)
    add_text(slide, content, x + Inches(0.15), y + Inches(0.65),
             Inches(5.8), Inches(1.85), font_size=9.5, color=C_DARK, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — MA TRẬN PHÂN QUYỀN
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('08. Ma trận Phân quyền theo Vai trò', accent=C_PRIMARY)

ROLES = ['CEO', 'Director', 'Mgr DH', 'Sale DH', 'Mgr NH', 'Sale NH', 'Ops', 'Finance']
MODULES_P = ['Dashboard','CRM','Booking DH','Booking NH','Proposal','Hợp đồng','Vận hành','Kế toán','Báo cáo']

perms = {
    'Dashboard':   [1,1,1,0,1,0,1,1],
    'CRM':         [0,1,1,1,1,1,0,0],
    'Booking DH':  [0,1,1,1,0,0,0,0],
    'Booking NH':  [0,1,1,0,1,1,0,0],
    'Proposal':    [1,1,1,0,0,0,0,0],
    'Hợp đồng':    [0,1,1,0,0,0,0,1],
    'Vận hành':    [0,1,0,0,0,0,1,0],
    'Kế toán':     [0,1,0,0,0,0,0,1],
    'Báo cáo':     [1,1,1,0,0,0,0,1],
}

col_w = Inches(1.4)
row_h = Inches(0.55)
header_x = Inches(1.6)
start_y = Inches(1.25)

# Header row — roles
for j, role in enumerate(ROLES):
    x = header_x + j * col_w
    add_rect(slide, x, start_y, col_w - Inches(0.05), row_h, fill=C_PRIMARY)
    add_text(slide, role, x, start_y + Pt(4), col_w - Inches(0.05), row_h - Pt(8),
             font_size=8.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Data rows
for i, mod in enumerate(MODULES_P):
    y = start_y + (i+1) * row_h
    bg = C_LIGHT_BG if i % 2 == 0 else C_WHITE
    add_rect(slide, Inches(0.35), y, Inches(1.2), row_h, fill=C_DIVIDER)
    add_text(slide, mod, Inches(0.4), y + Pt(4), Inches(1.15), row_h - Pt(8),
             font_size=9, bold=True, color=C_DARK)
    for j, role in enumerate(ROLES):
        x = header_x + j * col_w
        has_perm = perms.get(mod, [0]*8)[j]
        cell_color = RGBColor(0xD1,0xFA,0xE5) if has_perm else C_WHITE
        add_rect(slide, x, y, col_w - Inches(0.05), row_h - Pt(2),
                 fill=cell_color, line_color=C_DIVIDER, line_width=Pt(0.5))
        symbol = '✓' if has_perm else '·'
        sym_color = C_GREEN if has_perm else C_GRAY
        add_text(slide, symbol, x, y + Pt(4), col_w - Inches(0.05), row_h - Pt(8),
                 font_size=12, bold=has_perm, color=sym_color, align=PP_ALIGN.CENTER)

# Legend
add_rect(slide, Inches(0.35), Inches(7.0), Inches(3.5), Inches(0.28), fill=C_WHITE)
add_rect(slide, Inches(0.35), Inches(7.0), Inches(0.35), Inches(0.28), fill=C_GREEN)
add_text(slide, '✓ Có quyền truy cập', Inches(0.75), Inches(7.0), Inches(2.8), Inches(0.28),
         font_size=8.5, color=C_DARK)
add_text(slide, '  ·  Không có quyền', Inches(3.9), Inches(7.0), Inches(3), Inches(0.28),
         font_size=8.5, color=C_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — KẾT QUẢ TEST
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Kết quả Kiểm thử Hệ thống', accent=C_GREEN)

add_text(slide, 'Test toàn diện ngày 21/06/2026  —  16/19 kịch bản PASS',
         Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.5),
         font_size=14, bold=True, color=C_GREEN)

test_results_data = [
    ('✓', 'SCN-1', 'Login Sale Executive Dài hạn', 'le.minh.tuan@thiso.com', 'PASS', C_GREEN),
    ('✓', 'SCN-1', 'Tìm unit VACANT để booking', 'Unit L3-B02 (270m²)', 'PASS', C_GREEN),
    ('✓', 'SCN-2', 'Login Manager Dài hạn', 'nguyen.thu.hoa@thiso.com', 'PASS', C_GREEN),
    ('✓', 'SCN-3', 'Login Sale Executive Ngắn hạn', 'hoang.thi.mai@thiso.com', 'PASS', C_GREEN),
    ('✓', 'SCN-3', 'Tìm unit có Slot', 'Unit L3-B02 có 6 slots', 'PASS', C_GREEN),
    ('✓', 'SCN-4', 'Login Mall Director', 'do.thi.nga@thiso.com', 'PASS', C_GREEN),
    ('✓', 'SCN-4', 'Login Operation Staff', 'trinh.van.long@thiso.com', 'PASS', C_GREEN),
    ('✓', 'SCN-4', 'Xem Pending Approvals (Director)', '1 chờ duyệt', 'PASS', C_GREEN),
    ('✓', 'SCN-5', 'Login Kế Toán Trưởng', 'bui.thi.hoa@thiso.com', 'PASS', C_GREEN),
    ('✓', 'SCN-5', 'Xem danh sách hóa đơn', '5 hóa đơn', 'PASS', C_GREEN),
    ('✓', 'SCN-5', 'AR Aging Report', 'OK', 'PASS', C_GREEN),
    ('✓', 'SCN-5', 'Revenue Report', 'OK', 'PASS', C_GREEN),
    ('✓', 'SCN-6', 'Login CEO', 'ceo@thiso.com', 'PASS', C_GREEN),
    ('✓', 'SCN-6', 'CEO Dashboard', 'OK', 'PASS', C_GREEN),
    ('✓', 'SCN-6', 'Cross-Mall Dashboard', 'OK', 'PASS', C_GREEN),
    ('✓', 'SCN-6', 'CRM Stats', 'OK', 'PASS', C_GREEN),
    ('~', 'SCN-1', 'Tạo Lead (source enum)', 'Sửa sang BROKER → OK', 'WARN', RGBColor(0xD9,0x77,0x06)),
    ('~', 'SCN-3', 'Slot Booking conflict', 'Trùng SB-2026-00001, đổi ngày → OK', 'WARN', RGBColor(0xD9,0x77,0x06)),
    ('~', 'SCN-4', 'Tạo Ticket (tenantId)', 'Cần liên kết tenant, UI có sẵn', 'WARN', RGBColor(0xD9,0x77,0x06)),
]

for i, (icon, scn, desc, detail, status, color) in enumerate(test_results_data[:18]):
    row = i // 2; col = i % 2
    x = Inches(0.35) + col * Inches(6.45)
    y = Inches(1.85) + row * Inches(0.55)
    if row > 8: break
    add_rect(slide, x, y, Inches(6.1), Inches(0.5),
             fill=RGBColor(0xF0,0xFD,0xF4) if color == C_GREEN else RGBColor(0xFF,0xF7,0xED),
             line_color=C_DIVIDER, line_width=Pt(0.5))
    add_text(slide, icon, x + Inches(0.05), y + Pt(3), Inches(0.3), Inches(0.42),
             font_size=11, color=color, bold=True)
    add_text(slide, scn, x + Inches(0.35), y + Pt(4), Inches(0.6), Inches(0.38),
             font_size=7.5, bold=True, color=C_GRAY)
    add_text(slide, desc, x + Inches(0.95), y + Pt(4), Inches(3.2), Inches(0.38),
             font_size=8, color=C_DARK)
    add_text(slide, status, x + Inches(5.5), y + Pt(4), Inches(0.55), Inches(0.38),
             font_size=7.5, bold=True, color=color, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — BẢNG THÔNG TIN TÀI KHOẢN TỔNG HỢP
# ═══════════════════════════════════════════════════════════════════════════════
slide = content_slide('Tổng hợp Tài khoản — Thông tin Đăng nhập', accent=C_PRIMARY)

add_text(slide, 'URL: http://localhost:8080   |   Mật khẩu mặc định mới: Thiso@2024   |   Admin: Admin123!',
         Inches(0.35), Inches(1.2), Inches(12.62), Inches(0.4),
         font_size=11, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

full_accounts = [
    ('TỔNG GIÁM ĐỐC','CEO','ceo@thiso.com','User123!','Toàn bộ hệ thống','Executive', C_DARK),
    ('Giám Đốc Mall','MALL_DIRECTOR','do.thi.nga@thiso.com','Thiso@2024','Dashboard, Mặt bằng, CRM, Vận hành','Operations', C_PRIMARY),
    ('Trưởng Sale Dài Hạn','LEASING_MANAGER','nguyen.thu.hoa@thiso.com','Thiso@2024','CRM, Booking, Proposal, Hợp đồng','Leasing - DH', RGBColor(0xD9,0x77,0x06)),
    ('Sale Executive DH','LEASING_EXECUTIVE','le.minh.tuan@thiso.com','Thiso@2024','CRM, Booking dài hạn','Leasing - DH', RGBColor(0xD9,0x77,0x06)),
    ('Sale Executive DH','LEASING_EXECUTIVE','pham.thi.lan@thiso.com','Thiso@2024','CRM, Booking dài hạn','Leasing - DH', RGBColor(0xD9,0x77,0x06)),
    ('Trưởng Sale Ngắn Hạn','LEASING_MANAGER','tran.van.duc@thiso.com','Thiso@2024','CRM, Booking slot, Xác nhận','Leasing - NH', C_VIOLET),
    ('Sale Executive NH','LEASING_EXECUTIVE','hoang.thi.mai@thiso.com','Thiso@2024','Booking slot ngắn hạn','Leasing - NH', C_VIOLET),
    ('Sale Executive NH','LEASING_EXECUTIVE','vu.quoc.binh@thiso.com','Thiso@2024','Booking slot ngắn hạn','Leasing - NH', C_VIOLET),
    ('Vận Hành','OPERATION','trinh.van.long@thiso.com','Thiso@2024','Ticket, Fitout, Bảo trì','Operations', C_GREEN),
    ('Kế Toán Trưởng','FINANCE','bui.thi.hoa@thiso.com','Thiso@2024','Billing, Báo cáo, SAP','Finance', RGBColor(0x0F,0x76,0x6E)),
    ('Kế Toán Viên','FINANCE','dang.van.quang@thiso.com','Thiso@2024','Billing, Hóa đơn','Finance', RGBColor(0x0F,0x76,0x6E)),
]

headers = ['Vai trò', 'Role', 'Email đăng nhập', 'Mật khẩu', 'Phân hệ truy cập', 'Bộ phận']
col_widths = [Inches(1.8), Inches(1.55), Inches(2.85), Inches(1.1), Inches(3.75), Inches(1.3)]
x_starts = [Inches(0.18)]
for w in col_widths[:-1]:
    x_starts.append(x_starts[-1] + w + Inches(0.03))

# Header
for j, (hdr, w, xs) in enumerate(zip(headers, col_widths, x_starts)):
    add_rect(slide, xs, Inches(1.65), w, Inches(0.38), fill=C_PRIMARY)
    add_text(slide, hdr, xs + Inches(0.04), Inches(1.67), w - Inches(0.08), Inches(0.34),
             font_size=8, bold=True, color=C_WHITE)

# Rows
for i, (role, sys_role, email, pwd, modules_a, dept, color) in enumerate(full_accounts):
    y = Inches(2.05) + i * Inches(0.46)
    bg = C_WHITE if i % 2 == 0 else C_LIGHT_BG
    add_rect(slide, Inches(0.18), y, Inches(13.0), Inches(0.44), fill=bg)
    # Color dot
    add_rect(slide, Inches(0.18), y + Inches(0.05), Inches(0.08), Inches(0.34), fill=color)
    vals = [role, sys_role, email, pwd, modules_a, dept]
    for j, (val, w, xs) in enumerate(zip(vals, col_widths, x_starts)):
        add_text(slide, val, xs + Inches(0.04), y + Pt(3), w - Inches(0.08), Inches(0.38),
                 font_size=7.5 if j == 4 else 8, color=C_DARK if j != 1 else C_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=C_PRIMARY)
add_rect(slide, 0, Inches(2.8), W, Inches(1.9), fill=C_WHITE)
add_rect(slide, 0, Inches(2.8), Inches(0.2), Inches(1.9), fill=C_ACCENT)

add_text(slide, 'Liên hệ Hỗ trợ', Inches(0.5), Inches(2.95), Inches(12.33), Inches(0.55),
         font_size=28, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
support = [
    '📧  Admin hệ thống:  admin@thiso.com',
    '📱  Hotline IT:  0901 234 567',
    '🌐  Hệ thống:  http://localhost:8080',
    '📋  Tài liệu:  Liên hệ quản trị viên',
]
add_text(slide, '   |   '.join(support[:2]),
         Inches(0.5), Inches(3.55), Inches(12.33), Inches(0.45),
         font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, '   |   '.join(support[2:]),
         Inches(0.5), Inches(3.95), Inches(12.33), Inches(0.45),
         font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 'THISO LEASING PLATFORM',
         Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.6),
         font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'Nền tảng quản lý cho thuê thương mại toàn diện',
         Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.4),
         font_size=14, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)
add_rect(slide, 0, Inches(6.9), W, Inches(0.6), fill=C_ACCENT)
add_text(slide, 'Phiên bản 1.0  |  Tháng 6/2026  |  Bảo mật nội bộ — Không phân phối bên ngoài',
         Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.4),
         font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
output_path = r'd:\leasing-platform-thiso\THISO_Leasing_Platform_UserGuide.pptx'
prs.save(output_path)
print('PPT saved: ' + output_path)
print('Slides: ' + str(len(prs.slides)))
